# cmn/tools/tool/bedrock_converse_tools_pptx.py

import json
import logging
import os
import re
import tempfile
from dataclasses import dataclass, field

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

from cmn.tools.tool.bedrock_converse_tools_tool import AbstractBedrockConverseTool
from cmn.tools.tool.pptx import PptxChartBuilder

logger = logging.getLogger(__name__)

_CONFIG_DIR = os.path.join(
    os.path.dirname(__file__),
    "..", "config", "pptx",
)

################################################################################
# SECTION: BrandGuidelines
################################################################################

@dataclass
class BrandColors:
    primary:       str = "#1B3A6B"
    secondary:     str = "#2E86AB"
    accent:        str = "#F4A261"
    background:    str = "#FFFFFF"
    surface:       str = "#F5F7FA"
    text_dark:     str = "#1A1A2E"
    text_light:    str = "#FFFFFF"
    text_muted:    str = "#6B7280"
    success:       str = "#10B981"
    warning_color: str = "#F59E0B"
    danger:        str = "#EF4444"
    chart_palette: list = field(default_factory=lambda: [
        "#1B3A6B", "#2E86AB", "#F4A261",
        "#10B981", "#F59E0B", "#EF4444",
    ])

    def hex_to_rgb(self, hex_color: str) -> tuple:
        h = hex_color.lstrip("#")
        return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

    def get(self, name: str) -> str:
        return getattr(self, name, self.primary)


@dataclass
class BrandFonts:
    heading:         str = "Calibri"
    body:            str = "Calibri"
    mono:            str = "Courier New"
    heading_size:    int = 36
    subheading_size: int = 24
    body_size:       int = 18
    caption_size:    int = 14


@dataclass
class BrandSlide:
    width_inches:  float = 13.33
    height_inches: float = 7.5


@dataclass
class BrandRules:
    accent_bar_height_inches: float = 0.08
    accent_bar_position:      str   = "bottom"
    logo_position:            str   = "bottom_right"
    max_bullets_per_slide:    int   = 6
    max_words_per_bullet:     int   = 15
    cover_background:         str   = "primary"
    section_background:       str   = "secondary"
    content_background:       str   = "background"
    use_slide_numbers:        bool  = True
    tone:                     str   = "professional"


@dataclass
class BrandLayoutConfig:
    title_y:    float = 0.05
    subtitle_y: float = 0.55
    body_y:     float = 0.20
    quote_y:    float = 0.30
    attr_y:     float = 0.65
    split:      float = 0.50


@dataclass
class BrandGuidelines:
    brand_name: str         = "Default"
    colors:     BrandColors = field(default_factory=BrandColors)
    fonts:      BrandFonts  = field(default_factory=BrandFonts)
    slide:      BrandSlide  = field(default_factory=BrandSlide)
    rules:      BrandRules  = field(default_factory=BrandRules)
    layouts:    dict        = field(default_factory=dict)

    # ── Factories ─────────────────────────────────────────────────────────────

    @classmethod
    def from_json(cls, path: str) -> "BrandGuidelines":
        with open(path, "r") as f:
            data = json.load(f)

        def _filter(dc, raw: dict) -> dict:
            return {
                k: v for k, v in raw.items()
                if k in dc.__dataclass_fields__
            }

        return cls(
            brand_name = data.get("brand_name", "Default"),
            colors     = BrandColors(**_filter(BrandColors, data.get("colors", {}))),
            fonts      = BrandFonts( **_filter(BrandFonts,  data.get("fonts",  {}))),
            slide      = BrandSlide( **_filter(BrandSlide,  data.get("slide",  {}))),
            rules      = BrandRules( **_filter(BrandRules,  data.get("rules",  {}))),
            layouts    = data.get("layouts", {}),
        )

    @classmethod
    def load_all(cls, config_dir: str) -> dict:
        """
        Load every .json in config_dir.
        Returns { brand_name_lower: BrandGuidelines }.
        Falls back to { "default": BrandGuidelines() } if dir missing or empty.
        """
        brands = {}

        if not os.path.isdir(config_dir):
            logger.warning(
                "Brand config dir not found: %s — using built-in default",
                config_dir,
            )
            return {"default": cls()}

        for fname in sorted(os.listdir(config_dir)):
            if not fname.endswith(".json"):
                continue
            path = os.path.join(config_dir, fname)
            try:
                brand = cls.from_json(path)
                brands[brand.brand_name.lower()] = brand
                logger.info("Loaded brand '%s' from %s", brand.brand_name, fname)
            except Exception as exc:
                logger.warning("Failed to load brand file %s: %s", fname, exc)

        if not brands:
            logger.warning("No brand files loaded — using built-in default")
            brands["default"] = cls()

        return brands

    # ── Helpers ───────────────────────────────────────────────────────────────

    def layout(self, slide_type: str) -> BrandLayoutConfig:
        """Return BrandLayoutConfig for slide_type with safe defaults."""
        raw   = self.layouts.get(slide_type, {})
        valid = {
            k: v for k, v in raw.items()
            if k in BrandLayoutConfig.__dataclass_fields__
        }
        return BrandLayoutConfig(**valid)

    def content_chart_split(self) -> float:
        """
        Bullet/chart split ratio for content_chart slides.
        Read from layouts.content_chart.split in brand JSON.
        Defaults to 0.45 (45% bullets / 55% chart).
        """
        return self.layouts.get("content_chart", {}).get("split", 0.45)


################################################################################
# SECTION: Filename helper
################################################################################

def _safe_filename(title: str) -> str:
    """Convert title to a safe cross-platform filename."""
    safe = re.sub(r'[<>:"/\\|?*]', "-", title)
    safe = re.sub(r"[-\s]+", "_", safe).strip("_")
    return safe[:100] or "presentation"


################################################################################
# SECTION: PptxBedrockConverseTool
################################################################################

class PptxBedrockConverseTool(AbstractBedrockConverseTool):

    def __init__(self, config_dir: str = _CONFIG_DIR):
        name = "create_pptx"
        definition = {
            "toolSpec": {
                "name": name,
                "description": (
                    "Creates a branded PowerPoint presentation (.pptx) "
                    "with native editable charts. "
                    "Supports text slides and chart slides. "
                    "IMPORTANT: When the user asks for a chart, graph, trend, "
                    "or visual in the presentation, always use slide_type 'chart' "
                    "or 'content_chart'. "
                    "Do NOT put data into bullet points when a chart is more appropriate. "
                    "For chart slides: fetch data from the relevant tool first if not "
                    "already in conversation context, then reshape into chart_data. "
                    "The tool is domain-agnostic — works for sales, inventory, "
                    "stock prices, headcount, or any numeric data. "
                    "Brand guidelines (colors, fonts, layout) are applied automatically."
                ),
                "inputSchema": {
                    "json": {
                        "type": "object",
                        "properties": {
                            "title": {
                                "type":        "string",
                                "description": "Presentation title.",
                            },
                            "author": {
                                "type":        "string",
                                "description": "Author name for document properties.",
                            },
                            "brand": {
                                "type":        "string",
                                "description": (
                                    "Brand name to apply. "
                                    "Defaults to 'default' if omitted or not found."
                                ),
                                "default": "default",
                            },
                            "slides": {
                                "type":        "array",
                                "description": "Ordered list of slide definitions.",
                                "items": {
                                    "type": "object",
                                    "properties": {

                                        # ── Slide type ────────────────────────
                                        "slide_type": {
                                            "type": "string",
                                            "enum": [
                                                "cover",
                                                "section",
                                                "content",
                                                "two_column",
                                                "quote",
                                                "chart",
                                                "content_chart",
                                                "closing",
                                            ],
                                            "description": (
                                                "cover: title page. "
                                                "section: divider slide. "
                                                "content: title + bullets only — use for text insights with no data visual. "
                                                "two_column: side-by-side bullet lists. "
                                                "chart: title + full-width native chart. "
                                                "  USE THIS when showing trends, comparisons, or distributions. "
                                                "content_chart: bullets on left + chart on right. "
                                                "  USE THIS to show insights alongside the supporting data visual. "
                                                "  This is the PREFERRED layout when you have both analysis and data. "
                                                "quote: pull quote. "
                                                "closing: end slide."
                                            ),
                                        },

                                        # ── Common fields ─────────────────────
                                        "title": {
                                            "type":        "string",
                                            "description": "Slide title.",
                                        },
                                        "subtitle": {
                                            "type":        "string",
                                            "description": "Subtitle (cover and closing slides).",
                                        },

                                        # ── Bullet fields ─────────────────────
                                        "bullets": {
                                            "type":     "array",
                                            "items":    {"type": "string"},
                                            "maxItems": 6,
                                            "description": (
                                                "Bullet points for content and content_chart slides. "
                                                "Maximum 6. Each bullet under 15 words. "
                                                "Split across multiple slides if more are needed."
                                            ),
                                        },

                                        # ── Two-column fields ─────────────────
                                        "left_header": {
                                            "type":        "string",
                                            "description": "Left column header (two_column slides).",
                                        },
                                        "left_bullets": {
                                            "type":     "array",
                                            "items":    {"type": "string"},
                                            "maxItems": 6,
                                            "description": "Left column bullets (two_column slides).",
                                        },
                                        "right_header": {
                                            "type":        "string",
                                            "description": "Right column header (two_column slides).",
                                        },
                                        "right_bullets": {
                                            "type":     "array",
                                            "items":    {"type": "string"},
                                            "maxItems": 6,
                                            "description": "Right column bullets (two_column slides).",
                                        },

                                        # ── Quote fields ──────────────────────
                                        "quote": {
                                            "type":        "string",
                                            "description": "Pull quote text (quote slides).",
                                        },
                                        "attribution": {
                                            "type":        "string",
                                            "description": "Quote attribution (quote slides).",
                                        },

                                        # ── Chart fields ──────────────────────
                                        "chart_type": {
                                            "type": "string",
                                            "enum": ["bar", "line", "pie"],
                                            "description": (
                                                "Chart type. REQUIRED for chart and content_chart slides. "
                                                "bar: comparisons between categories. "
                                                "line: trends over time — USE THIS for monthly/time-series data. "
                                                "pie: proportions (max 6 slices recommended)."
                                            ),
                                        },
                                        "chart_data": {
                                            "type": "array",
                                            "description": (
                                                "REQUIRED for chart and content_chart slides. "
                                                "Fetch from the relevant data tool first if not in "
                                                "conversation context, then reshape into this format. "
                                                "The tool is domain-agnostic — works for any numeric data. "
                                                "For monthly sales trend: use month_name as label, revenue as value. "
                                                "For comparisons: use category/product/region as label. "
                                                "For multi-year: add series field with year as value. "
                                                "Single series:  [{\"label\": \"January\", \"value\": 213000}, ...]. "
                                                "Multi-series:   [{\"label\": \"January\", \"value\": 213000, \"series\": \"2024\"}, ...]."
                                            ),
                                            "items": {
                                                "type": "object",
                                                "properties": {
                                                    "label": {
                                                        "type":        "string",
                                                        "description": (
                                                            "Category label shown on x-axis. "
                                                            "Examples: month name, product name, "
                                                            "region, date, department."
                                                        ),
                                                    },
                                                    "value": {
                                                        "type":        "number",
                                                        "description": (
                                                            "Numeric value for this data point. "
                                                            "Examples: revenue, units, price, count, percentage."
                                                        ),
                                                    },
                                                    "series": {
                                                        "type":        "string",
                                                        "description": (
                                                            "Series name for multi-series charts. "
                                                            "Examples: year (2023, 2024), region (North, South). "
                                                            "Omit for single-series charts."
                                                        ),
                                                    },
                                                },
                                                "required": ["label", "value"],
                                            },
                                        },
                                        "chart_title": {
                                            "type":        "string",
                                            "description": (
                                                "Title shown inside the chart. "
                                                "Defaults to slide title if omitted."
                                            ),
                                        },
                                        "x_label": {
                                            "type":        "string",
                                            "description": "X-axis label. Example: 'Month', 'Product', 'Region'.",
                                        },
                                        "y_label": {
                                            "type":        "string",
                                            "description": "Y-axis label. Example: 'Revenue ($)', 'Units Sold'.",
                                        },

                                        # ── Speaker notes ─────────────────────
                                        "speaker_notes": {
                                            "type":        "string",
                                            "description": "Speaker notes for this slide.",
                                        },
                                    },
                                    "required": ["slide_type"],
                                },
                            },
                        },
                        "required": ["slides"],
                    }
                },
            }
        }
        super().__init__(name, definition)

        self._brands        = BrandGuidelines.load_all(config_dir)
        self._chart_builder = PptxChartBuilder()

    # ── Summary ───────────────────────────────────────────────────────────────

    def summary(self) -> str:
        brands = ", ".join(sorted(self._brands.keys()))
        return (
            "create_pptx : creates a branded PowerPoint (.pptx) "
            "with native editable charts. "

            # ── Slide types ───────────────────────────────────────────────────────
            "Slide types: "
            "cover (title page), "
            "section (divider), "
            "content (title + bullets — text only, NO numeric data), "
            "two_column (side-by-side bullets), "
            "chart (title + full-width native chart), "
            "content_chart (bullets on left + chart on right — "
            "PREFERRED when you have both insights and data), "
            "quote (pull quote), "
            "closing (end slide). "

            # ── Hard rule ─────────────────────────────────────────────────────────
            "RULE: When you have numeric data (revenue, units, prices, counts) "
            "always use slide_type 'chart' or 'content_chart'. "
            "NEVER put time-series or comparative numeric data into bullet points only. "

            # ── Chart types ───────────────────────────────────────────────────────
            "Chart types: "
            "bar (comparisons between categories), "
            "line (trends over time — use for monthly/time-series data), "
            "pie (proportions). "

            # ── Concrete reshaping example ────────────────────────────────────────
            "HOW TO BUILD A CHART SLIDE from sales_data result: "
            "sales_data returns rows like "
            "[{month_name: 'January', revenue: 213000, units_sold: 1035}, ...]. "
            "Reshape into a content_chart slide like this: "
            "{ "
            "slide_type: 'content_chart', "
            "title: 'Monthly Revenue 2024', "
            "chart_type: 'line', "
            "x_label: 'Month', "
            "y_label: 'Revenue ($)', "
            "chart_data: [ "
            "{label: 'January', value: 213000}, "
            "{label: 'February', value: 198000}, "
            "... one entry per month ... "
            "], "
            "bullets: ['Total $3.13M', 'Peak Dec $415K', 'Strong H2 growth'] "
            "} "
            "For YoY comparison use series field: "
            "{label: 'January', value: 213000, series: '2024'}. "

            # ── Recommended structure ─────────────────────────────────────────────
            "RECOMMENDED slide structure for a data presentation: "
            "1. cover, "
            "2. content (executive summary bullets), "
            "3. content_chart (primary metric line/bar chart + key insights), "
            "4. content_chart (secondary metric chart + insights), "
            "5. two_column or content (qualitative analysis), "
            "6. closing. "
            "ALWAYS include at least one content_chart or chart slide "
            "when numeric data is available. "

            f"Available brands: {brands}."
        )

    # ── Invoke ────────────────────────────────────────────────────────────────

    def invoke(self, params, tool_args: dict = None) -> dict:
        args      = tool_args or {}
        slides    = args.get("slides",  [])
        title     = args.get("title",   "Presentation")
        author    = args.get("author",  "")
        brand_key = args.get("brand",   "default").lower()

        if not slides:
            return {
                "error": (
                    "No slides provided. "
                    "Pass a list of slide definitions with slide_type, title, "
                    "bullets, and/or chart_data."
                )
            }

        # ── Resolve brand ─────────────────────────────────────────────────────
        warnings = []
        brand    = self._brands.get(brand_key)

        if brand is None:
            warnings.append(
                f"Brand '{brand_key}' not found. "
                f"Available: {list(self._brands.keys())}. Using default."
            )
            brand = self._brands.get("default", BrandGuidelines())

        # ── Build presentation ────────────────────────────────────────────────
        try:
            prs = self._build_presentation(brand, slides, title, author, warnings)
        except Exception as exc:
            logger.exception("PptxTool: build failed")
            return {"error": f"Failed to build presentation: {exc}"}

        # ── Save to temp file ─────────────────────────────────────────────────
        safe_title = _safe_filename(title)
        filename   = f"{safe_title}.pptx"

        tmp = tempfile.NamedTemporaryFile(
            suffix  = ".pptx",
            prefix  = f"pptx_{safe_title}_",
            delete  = False,
        )
        prs.save(tmp.name)
        tmp.close()

        logger.info("PptxTool: saved '%s' → %s", filename, tmp.name)

        return {
            "status":       "pptx_ready",
            "title":        title,
            "brand":        brand.brand_name,
            "slide_count":  len(slides),
            "slide_titles": [s.get("title", "") for s in slides],
            "filepath":     tmp.name,
            "filename":     filename,
            "warnings":     warnings,
        }

    # ── Presentation Builder ──────────────────────────────────────────────────

    def _build_presentation(
        self,
        brand:    BrandGuidelines,
        slides:   list,
        title:    str,
        author:   str,
        warnings: list,
    ) -> Presentation:

        prs              = Presentation()
        prs.slide_width  = Inches(brand.slide.width_inches)
        prs.slide_height = Inches(brand.slide.height_inches)

        core        = prs.core_properties
        core.title  = title
        core.author = author

        blank_layout = prs.slide_layouts[6]   # fully blank

        for idx, slide_def in enumerate(slides):
            slide_type = slide_def.get("slide_type", "content")
            slide      = prs.slides.add_slide(blank_layout)

            builder = self._BUILDERS.get(slide_type, self._build_content)
            builder(self, slide, slide_def, brand, warnings)

            self._add_accent_bar(slide, brand)

            if brand.rules.use_slide_numbers:
                self._add_slide_number(slide, brand, idx + 1)

            notes = slide_def.get("speaker_notes", "")
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

        return prs

    # ── Slide Builders ────────────────────────────────────────────────────────

    def _build_cover(self, slide, slide_def, brand, warnings):
        W  = brand.slide.width_inches
        H  = brand.slide.height_inches
        lc = brand.layout("cover")

        self._fill_background(
            slide,
            brand.colors.hex_to_rgb(brand.colors.get(brand.rules.cover_background)),
        )

        if title := slide_def.get("title", ""):
            self._add_text_box(
                slide,
                text      = title,
                left      = Inches(0.8),
                top       = Inches(H * lc.title_y),
                width     = Inches(W - 1.6),
                height    = Inches(1.2),
                font_name = brand.fonts.heading,
                font_size = brand.fonts.heading_size + 8,
                bold      = True,
                color_rgb = brand.colors.hex_to_rgb(brand.colors.text_light),
                align     = PP_ALIGN.LEFT,
            )

        if subtitle := slide_def.get("subtitle", ""):
            self._add_text_box(
                slide,
                text      = subtitle,
                left      = Inches(0.8),
                top       = Inches(H * lc.subtitle_y),
                width     = Inches(W - 1.6),
                height    = Inches(0.8),
                font_name = brand.fonts.body,
                font_size = brand.fonts.subheading_size,
                bold      = False,
                color_rgb = brand.colors.hex_to_rgb(brand.colors.accent),
                align     = PP_ALIGN.LEFT,
            )

    def _build_section(self, slide, slide_def, brand, warnings):
        W  = brand.slide.width_inches
        H  = brand.slide.height_inches
        lc = brand.layout("section")

        self._fill_background(
            slide,
            brand.colors.hex_to_rgb(brand.colors.get(brand.rules.section_background)),
        )

        if title := slide_def.get("title", ""):
            self._add_text_box(
                slide,
                text      = title,
                left      = Inches(0.8),
                top       = Inches(H * lc.title_y),
                width     = Inches(W - 1.6),
                height    = Inches(1.2),
                font_name = brand.fonts.heading,
                font_size = brand.fonts.heading_size + 4,
                bold      = True,
                color_rgb = brand.colors.hex_to_rgb(brand.colors.text_light),
                align     = PP_ALIGN.CENTER,
            )

    def _build_content(self, slide, slide_def, brand, warnings):
        W  = brand.slide.width_inches
        H  = brand.slide.height_inches
        lc = brand.layout("content")

        self._fill_background(
            slide,
            brand.colors.hex_to_rgb(brand.colors.get(brand.rules.content_background)),
        )
        self._add_title_bar(slide, brand)

        if title := slide_def.get("title", ""):
            self._add_text_box(
                slide,
                text      = title,
                left      = Inches(0.4),
                top       = Inches(H * lc.title_y),
                width     = Inches(W - 0.8),
                height    = Inches(0.6),
                font_name = brand.fonts.heading,
                font_size = brand.fonts.subheading_size,
                bold      = True,
                color_rgb = brand.colors.hex_to_rgb(brand.colors.primary),
                align     = PP_ALIGN.LEFT,
            )

        bullets = slide_def.get("bullets", [])
        if len(bullets) > brand.rules.max_bullets_per_slide:
            warnings.append(
                f"Slide '{slide_def.get('title', '')}': "
                f"{len(bullets)} bullets exceeds recommended max "
                f"{brand.rules.max_bullets_per_slide}. "
                "Consider splitting into two slides."
            )

        if bullets:
            self._add_bullet_box(
                slide,
                bullets = bullets,
                left    = Inches(0.5),
                top     = Inches(H * lc.body_y),
                width   = Inches(W - 1.0),
                height  = Inches(H - lc.body_y - 0.5),
                brand   = brand,
            )

    def _build_two_column(self, slide, slide_def, brand, warnings):
        W  = brand.slide.width_inches
        H  = brand.slide.height_inches
        lc = brand.layout("two_column")

        self._fill_background(
            slide,
            brand.colors.hex_to_rgb(brand.colors.get(brand.rules.content_background)),
        )
        self._add_title_bar(slide, brand)

        if title := slide_def.get("title", ""):
            self._add_text_box(
                slide,
                text      = title,
                left      = Inches(0.4),
                top       = Inches(H * lc.title_y),
                width     = Inches(W - 0.8),
                height    = Inches(0.6),
                font_name = brand.fonts.heading,
                font_size = brand.fonts.subheading_size,
                bold      = True,
                color_rgb = brand.colors.hex_to_rgb(brand.colors.primary),
                align     = PP_ALIGN.LEFT,
            )

        col_w    = (W - 1.2) * lc.split
        body_top = Inches(H * lc.body_y)
        body_h   = Inches(H - lc.body_y - 0.5)

        for side, x_left, header_key, bullets_key in [
            ("left",  Inches(0.5),         "left_header",  "left_bullets"),
            ("right", Inches(0.7 + col_w), "right_header", "right_bullets"),
        ]:
            if header := slide_def.get(header_key, ""):
                self._add_text_box(
                    slide,
                    text      = header,
                    left      = x_left,
                    top       = body_top,
                    width     = Inches(col_w),
                    height    = Inches(0.4),
                    font_name = brand.fonts.heading,
                    font_size = brand.fonts.body_size,
                    bold      = True,
                    color_rgb = brand.colors.hex_to_rgb(brand.colors.secondary),
                    align     = PP_ALIGN.LEFT,
                )
            if bullets := slide_def.get(bullets_key, []):
                if len(bullets) > brand.rules.max_bullets_per_slide:
                    warnings.append(
                        f"Slide '{slide_def.get('title', '')}' {side} column: "
                        f"{len(bullets)} bullets exceeds recommended max "
                        f"{brand.rules.max_bullets_per_slide}."
                    )
                self._add_bullet_box(
                    slide,
                    bullets = bullets,
                    left    = x_left,
                    top     = Inches(H * lc.body_y + 0.45),
                    width   = Inches(col_w),
                    height  = body_h,
                    brand   = brand,
                )

        self._add_vertical_divider(slide, brand, col_w, lc.body_y, H)

    def _build_quote(self, slide, slide_def, brand, warnings):
        W  = brand.slide.width_inches
        H  = brand.slide.height_inches
        lc = brand.layout("quote")

        self._fill_background(
            slide,
            brand.colors.hex_to_rgb(brand.colors.surface),
        )

        if quote := slide_def.get("quote", ""):
            self._add_text_box(
                slide,
                text      = f"\u201c{quote}\u201d",
                left      = Inches(1.2),
                top       = Inches(H * lc.quote_y),
                width     = Inches(W - 2.4),
                height    = Inches(2.0),
                font_name = brand.fonts.heading,
                font_size = brand.fonts.heading_size - 4,
                bold      = False,
                italic    = True,
                color_rgb = brand.colors.hex_to_rgb(brand.colors.primary),
                align     = PP_ALIGN.CENTER,
            )

        if attribution := slide_def.get("attribution", ""):
            self._add_text_box(
                slide,
                text      = attribution,
                left      = Inches(1.2),
                top       = Inches(H * lc.attr_y),
                width     = Inches(W - 2.4),
                height    = Inches(0.4),
                font_name = brand.fonts.body,
                font_size = brand.fonts.caption_size,
                bold      = False,
                color_rgb = brand.colors.hex_to_rgb(brand.colors.text_muted),
                align     = PP_ALIGN.RIGHT,
            )

    def _build_chart(self, slide, slide_def, brand, warnings):
        """Full-width native chart slide."""
        W  = brand.slide.width_inches
        H  = brand.slide.height_inches
        lc = brand.layout("content")

        self._fill_background(
            slide,
            brand.colors.hex_to_rgb(brand.colors.get(brand.rules.content_background)),
        )
        self._add_title_bar(slide, brand)

        if title := slide_def.get("title", ""):
            self._add_text_box(
                slide,
                text      = title,
                left      = Inches(0.4),
                top       = Inches(H * lc.title_y),
                width     = Inches(W - 0.8),
                height    = Inches(0.6),
                font_name = brand.fonts.heading,
                font_size = brand.fonts.subheading_size,
                bold      = True,
                color_rgb = brand.colors.hex_to_rgb(brand.colors.primary),
                align     = PP_ALIGN.LEFT,
            )

        self._chart_builder.add_chart(
            slide     = slide,
            slide_def = slide_def,
            brand     = brand,
            warnings  = warnings,
            left      = Inches(0.5),
            top       = Inches(H * lc.body_y),
            width     = Inches(W - 1.0),
            height    = Inches(H - lc.body_y - 0.4),
        )

    def _build_content_chart(self, slide, slide_def, brand, warnings):
        """Left: bullets  |  Right: native chart."""
        W  = brand.slide.width_inches
        H  = brand.slide.height_inches
        lc = brand.layout("content")

        self._fill_background(
            slide,
            brand.colors.hex_to_rgb(brand.colors.get(brand.rules.content_background)),
        )
        self._add_title_bar(slide, brand)

        if title := slide_def.get("title", ""):
            self._add_text_box(
                slide,
                text      = title,
                left      = Inches(0.4),
                top       = Inches(H * lc.title_y),
                width     = Inches(W - 0.8),
                height    = Inches(0.6),
                font_name = brand.fonts.heading,
                font_size = brand.fonts.subheading_size,
                bold      = True,
                color_rgb = brand.colors.hex_to_rgb(brand.colors.primary),
                align     = PP_ALIGN.LEFT,
            )

        # ── Split ratio from brand JSON, default 45/55 ────────────────────────
        split    = brand.content_chart_split()
        col_w    = (W - 1.2) * split
        body_top = H * lc.body_y

        # ── Left: bullets ─────────────────────────────────────────────────────
        if bullets := slide_def.get("bullets", []):
            if len(bullets) > brand.rules.max_bullets_per_slide:
                warnings.append(
                    f"Slide '{slide_def.get('title', '')}': "
                    f"{len(bullets)} bullets exceeds recommended max "
                    f"{brand.rules.max_bullets_per_slide}."
                )
            self._add_bullet_box(
                slide,
                bullets = bullets,
                left    = Inches(0.5),
                top     = Inches(body_top),
                width   = Inches(col_w),
                height  = Inches(H - body_top - 0.5),
                brand   = brand,
            )

        # ── Divider ───────────────────────────────────────────────────────────
        self._add_vertical_divider(slide, brand, col_w, lc.body_y, H)

        # ── Right: native chart ───────────────────────────────────────────────
        chart_left = 0.7 + col_w
        self._chart_builder.add_chart(
            slide     = slide,
            slide_def = slide_def,
            brand     = brand,
            warnings  = warnings,
            left      = Inches(chart_left),
            top       = Inches(body_top),
            width     = Inches(W - chart_left - 0.3),
            height    = Inches(H - body_top - 0.4),
        )

    def _build_closing(self, slide, slide_def, brand, warnings):
        """Closing reuses cover layout."""
        self._build_cover(slide, slide_def, brand, warnings)

    # ── Dispatch table ────────────────────────────────────────────────────────

    _BUILDERS = {
        "cover":         _build_cover,
        "section":       _build_section,
        "content":       _build_content,
        "two_column":    _build_two_column,
        "quote":         _build_quote,
        "chart":         _build_chart,
        "content_chart": _build_content_chart,
        "closing":       _build_closing,
    }

    # ── Drawing Helpers ───────────────────────────────────────────────────────

    def _fill_background(self, slide, rgb: tuple):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*rgb)

    # bedrock_converse_tools_pptx.py

    def _add_title_bar(self, slide, brand: BrandGuidelines):
        """Thin primary-color rule under the title area."""
        W = brand.slide.width_inches

        # ── Layout ────────────────────────────────────────────────────────────────
        # title text:  top=0.375"  height=0.6"  bottom=0.975"
        # upper space: 0.15"  gap between title bottom and bar top
        # bar top:     0.975 + 0.15 = 1.125"
        # bar height:  0.03"
        # lower space: handled by body_y in brand JSON

        shape = slide.shapes.add_shape(
            1,
            Inches(0.4),       Inches(1.125),    # ← was 1.0, now 1.125
            Inches(W - 0.8),   Inches(0.03),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(
            *brand.colors.hex_to_rgb(brand.colors.primary)
        )
        shape.line.fill.background()

    def _add_accent_bar(self, slide, brand: BrandGuidelines):
        """Full-width accent bar at bottom (or top) of every slide."""
        W     = brand.slide.width_inches
        H     = brand.slide.height_inches
        bar_h = brand.rules.accent_bar_height_inches
        top   = (
            Inches(H - bar_h)
            if brand.rules.accent_bar_position == "bottom"
            else Inches(0)
        )
        shape = slide.shapes.add_shape(
            1,
            Inches(0), top, Inches(W), Inches(bar_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(
            *brand.colors.hex_to_rgb(brand.colors.accent)
        )
        shape.line.fill.background()

    def _add_slide_number(self, slide, brand: BrandGuidelines, number: int):
        W = brand.slide.width_inches
        H = brand.slide.height_inches
        self._add_text_box(
            slide,
            text      = str(number),
            left      = Inches(W - 0.5),
            top       = Inches(H - 0.35),
            width     = Inches(0.4),
            height    = Inches(0.25),
            font_name = brand.fonts.body,
            font_size = brand.fonts.caption_size - 2,
            bold      = False,
            color_rgb = brand.colors.hex_to_rgb(brand.colors.text_muted),
            align     = PP_ALIGN.RIGHT,
        )

    def _add_vertical_divider(
        self,
        slide,
        brand:  BrandGuidelines,
        col_w:  float,
        body_y: float,
        H:      float,
    ):
        shape = slide.shapes.add_shape(
            1,
            Inches(0.6 + col_w), Inches(H * body_y),
            Inches(0.02),        Inches(H - body_y - 0.4),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(
            *brand.colors.hex_to_rgb(brand.colors.secondary)
        )
        shape.line.fill.background()

    def _add_text_box(
        self,
        slide,
        text:      str,
        left:      Emu,
        top:       Emu,
        width:     Emu,
        height:    Emu,
        font_name: str,
        font_size: int,
        bold:      bool,
        color_rgb: tuple,
        align:     PP_ALIGN = PP_ALIGN.LEFT,
        italic:    bool = False,
    ):
        tf           = slide.shapes.add_textbox(left, top, width, height).text_frame
        tf.word_wrap = True
        p            = tf.paragraphs[0]
        p.alignment  = align
        run          = p.add_run()
        run.text           = text
        run.font.name      = font_name
        run.font.size      = Pt(font_size)
        run.font.bold      = bold
        run.font.italic    = italic
        run.font.color.rgb = RGBColor(*color_rgb)

    def _add_bullet_box(
        self,
        slide,
        bullets: list,
        left:    Emu,
        top:     Emu,
        width:   Emu,
        height:  Emu,
        brand:   BrandGuidelines,
    ):
        tf           = slide.shapes.add_textbox(left, top, width, height).text_frame
        tf.word_wrap = True
        text_rgb     = brand.colors.hex_to_rgb(brand.colors.text_dark)
        accent_rgb   = brand.colors.hex_to_rgb(brand.colors.accent)

        for i, bullet in enumerate(bullets):
            p              = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment    = PP_ALIGN.LEFT
            p.space_before = Pt(4)

            # Bullet marker in accent color
            marker                = p.add_run()
            marker.text           = "▪  "
            marker.font.name      = brand.fonts.body
            marker.font.size      = Pt(brand.fonts.body_size)
            marker.font.color.rgb = RGBColor(*accent_rgb)

            # Bullet text in dark color
            run                = p.add_run()
            run.text           = bullet
            run.font.name      = brand.fonts.body
            run.font.size      = Pt(brand.fonts.body_size)
            run.font.color.rgb = RGBColor(*text_rgb)